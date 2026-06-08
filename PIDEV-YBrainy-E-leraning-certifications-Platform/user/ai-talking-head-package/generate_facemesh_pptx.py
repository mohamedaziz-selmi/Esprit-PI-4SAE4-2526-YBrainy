from __future__ import annotations

from pathlib import Path

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from scipy.stats import zscore
from sklearn.decomposition import PCA
from sklearn.preprocessing import StandardScaler


ROOT = Path(__file__).resolve().parent
OUTPUT_DIR = ROOT / "report_outputs"
ASSETS_DIR = OUTPUT_DIR / "assets"
PPTX_PATH = OUTPUT_DIR / "FaceMesh_Feature_Extraction_Report_Updated.pptx"

BG = RGBColor(247, 245, 240)
NAVY = RGBColor(22, 36, 71)
TEAL = RGBColor(41, 110, 118)
CORAL = RGBColor(218, 110, 86)
OLIVE = RGBColor(117, 135, 85)
INK = RGBColor(38, 43, 54)
MUTED = RGBColor(94, 100, 114)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(232, 236, 241)

FEATURE_COLS = [
    "eye_left_open",
    "eye_right_open",
    "nose_width_ratio",
    "nose_bridge_ratio",
    "face_shape_index",
    "jaw_width_ratio",
    "lip_fullness_ratio",
    "eyebrow_arch_ratio",
    "interpupillary_ratio",
    "chin_projection_ratio",
]
FEATURE_LABELS = {
    "eye_left_open": "Ouverture oeil gauche",
    "eye_right_open": "Ouverture oeil droit",
    "nose_width_ratio": "Largeur du nez",
    "nose_bridge_ratio": "Longueur de l'arete nasale",
    "face_shape_index": "Indice de forme du visage",
    "jaw_width_ratio": "Largeur de la machoire",
    "lip_fullness_ratio": "Epaisseur des levres",
    "eyebrow_arch_ratio": "Arc du sourcil",
    "interpupillary_ratio": "Distance interpupillaire",
    "chin_projection_ratio": "Projection du menton",
}


def load_metrics() -> dict[str, object]:
    features = pd.read_csv(OUTPUT_DIR / "sample_features.csv")
    best_models = pd.read_csv(OUTPUT_DIR / "best_model_per_family.csv")
    cluster_profiles = pd.read_csv(OUTPUT_DIR / "cluster_profiles.csv")
    detection_by_race = pd.read_csv(OUTPUT_DIR / "detection_rate_by_race.csv")

    detected = features[features["detected"] == True].copy()
    mask = (np.abs(zscore(detected[FEATURE_COLS], nan_policy="omit")) < 3).all(axis=1)
    clean = detected.loc[mask, FEATURE_COLS].copy()
    scaler = StandardScaler()
    X = scaler.fit_transform(clean)
    pca = PCA(n_components=0.85)
    pca.fit(X)

    loadings = pd.DataFrame(
        pca.components_.T,
        index=FEATURE_COLS,
        columns=[f"PC{i + 1}" for i in range(pca.n_components_)],
    )

    top_features = []
    for component in loadings.columns[:4]:
        names = (
            loadings[component]
            .abs()
            .sort_values(ascending=False)
            .head(3)
            .index.map(FEATURE_LABELS.get)
            .tolist()
        )
        top_features.append((component, names))

    return {
        "features": features,
        "best_models": best_models,
        "cluster_profiles": cluster_profiles,
        "detection_by_race": detection_by_race,
        "sample_size": int(len(features)),
        "detected_faces": int(features["detected"].sum()),
        "detection_rate": float(features["detected"].mean()),
        "duplicates": int(features["image_md5"].dropna().duplicated().sum()),
        "outliers_removed": int((~mask).sum()),
        "modeling_rows": int(mask.sum()),
        "pca_components": int(pca.n_components_),
        "pca_variance": float(pca.explained_variance_ratio_.sum()),
        "pca_ratios": pca.explained_variance_ratio_,
        "top_features": top_features,
    }


def add_slide(prs: Presentation) -> object:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.7), Inches(0.8))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = NAVY
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(1.03), Inches(11.5), Inches(0.45))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.name = "Aptos"
        srun.font.size = Pt(11)
        srun.font.color.rgb = MUTED


def add_accent_bar(slide) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.22))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, font_size=18, color=INK, bold=False, name="Aptos") -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_bullets(slide, left, top, width, height, bullets: list[str], color=INK, font_size=16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def add_card(slide, left, top, width, height, title, value, fill_color) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Aptos"
    r1.font.size = Pt(12)
    r1.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = value
    r2.font.name = "Aptos Display"
    r2.font.size = Pt(22)
    r2.font.bold = True
    r2.font.color.rgb = WHITE


def add_table(slide, left, top, width, height, headers: list[str], rows: list[list[str]]) -> None:
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else LIGHT
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(10)
                    run.font.color.rgb = INK


def add_picture(slide, image_name: str, left, top, width=None, height=None) -> None:
    path = ASSETS_DIR / image_name
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def build_presentation() -> Path:
    metrics = load_metrics()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    best_model = metrics["best_models"].iloc[0]
    kmeans_model = metrics["best_models"][metrics["best_models"]["model"] == "KMeans"].iloc[0]
    detection_min = metrics["detection_by_race"]["rate"].min()
    detection_max = metrics["detection_by_race"]["rate"].max()

    slide = add_slide(prs)
    add_accent_bar(slide)
    add_textbox(slide, Inches(0.58), Inches(0.78), Inches(7.4), Inches(0.55), "FaceMesh Feature Extraction", 28, NAVY, True, "Aptos Display")
    add_textbox(slide, Inches(0.6), Inches(1.46), Inches(6.6), Inches(1.1), "FairFace, ACP et clustering pour l'analyse des traits faciaux", 17, TEAL, False)
    add_bullets(
        slide,
        Inches(0.65),
        Inches(2.4),
        Inches(5.4),
        Inches(2.2),
        [
            "Rapport de synthese en format presentation",
            "Pipeline execute sur un echantillon stratifie de 10 000 images",
            "Sorties coherentes avec le document .docx mis a jour",
        ],
        INK,
        16,
    )
    add_card(slide, Inches(7.85), Inches(0.95), Inches(1.65), Inches(1.35), "Echantillon", f"{metrics['sample_size']:,}".replace(",", " "), TEAL)
    add_card(slide, Inches(9.68), Inches(0.95), Inches(1.65), Inches(1.35), "Detection", f"{metrics['detection_rate'] * 100:.2f}%", CORAL)
    add_card(slide, Inches(11.5), Inches(0.95), Inches(1.35), Inches(1.35), "ACP", f"{metrics['pca_components']} PC", OLIVE)
    add_picture(slide, "cluster_scatter.png", Inches(6.9), Inches(2.05), width=Inches(5.9))

    slide = add_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "1. Executive Summary", "Resultats cles retenus pour le rendu final.")
    add_card(slide, Inches(0.7), Inches(1.5), Inches(2.25), Inches(1.25), "Visages detectes", str(metrics["detected_faces"]), NAVY)
    add_card(slide, Inches(3.15), Inches(1.5), Inches(2.25), Inches(1.25), "Lignes modelisees", str(metrics["modeling_rows"]), TEAL)
    add_card(slide, Inches(5.6), Inches(1.5), Inches(2.25), Inches(1.25), "Outliers retires", str(metrics["outliers_removed"]), CORAL)
    add_card(slide, Inches(8.05), Inches(1.5), Inches(2.45), Inches(1.25), "Meilleur modele", f"{best_model['model_label']} / k={int(best_model['k'])}", OLIVE)
    add_bullets(
        slide,
        Inches(0.8),
        Inches(3.15),
        Inches(5.9),
        Inches(2.6),
        [
            "Le pipeline transforme les images FairFace en 10 ratios geometriques interpretable.",
            f"L'ACP conserve {metrics['pca_variance'] * 100:.2f}% de la variance avec {metrics['pca_components']} composantes.",
            f"Le meilleur score de silhouette observe est {best_model['silhouette']:.3f} avec {best_model['model_label']}.",
            f"KMeans reste une option compatible avec le pipeline voix precedent (silhouette {kmeans_model['silhouette']:.3f}).",
        ],
    )
    add_picture(slide, "model_comparison.png", Inches(7.35), Inches(3.0), width=Inches(5.35))

    slide = add_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "2. Business Understanding", "Business Objectives et Data Science Objectives consolides.")
    add_table(
        slide,
        Inches(0.6),
        Inches(1.45),
        Inches(12.1),
        Inches(4.35),
        ["ID", "Business Objective", "Data Science Objective", "Indicateur"],
        [
            ["BO-1", "Automatiser l'extraction de traits faciaux sur CPU", "Detecter le visage et calculer 10 ratios geometriques via FaceMesh", f"{metrics['detection_rate'] * 100:.2f}% de detection"],
            ["BO-2", "Obtenir une representation compacte pour le Pi", "Standardiser puis appliquer une ACP a 85%+ de variance", f"{metrics['pca_components']} composantes / {metrics['pca_variance'] * 100:.2f}%"],
            ["BO-3", "Segmenter les profils faciaux", "Comparer plusieurs modeles non supervises et retenir le meilleur compromis", f"{best_model['model_label']} / silhouette {best_model['silhouette']:.3f}"],
            ["BO-4", "Suivre la couverture des sous-groupes", "Mesurer la detection par race et documenter les ecarts", f"Entre {detection_min:.1f}% et {detection_max:.1f}%"],
        ],
    )
    add_textbox(slide, Inches(0.75), Inches(6.15), Inches(11.7), Inches(0.5), "Point de cadrage : les clusters sont des regroupements geometriques utiles a l'indexation et a la personnalisation, pas des categories biologiques fixes.", 12, MUTED)

    slide = add_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "3. Data Understanding", "Vue d'ensemble du dataset FairFace et EDA initiale.")
    add_picture(slide, "race_distribution.png", Inches(0.65), Inches(1.5), width=Inches(6.0))
    add_picture(slide, "age_distribution.png", Inches(6.95), Inches(1.5), width=Inches(5.7))
    add_bullets(
        slide,
        Inches(0.8),
        Inches(5.0),
        Inches(12.0),
        Inches(1.6),
        [
            "Le dataset contient 97 698 images annotees et reste relativement equilibre sur 7 groupes raciaux.",
            "Les tranches 20-29 et 30-39 sont les plus representees, ce qui doit etre garde en tete lors de l'interpretation.",
            "La repartition du genre reste proche d'un partage 50/50, favorable a une EDA defendable.",
        ],
        font_size=15,
    )

    slide = add_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "4. EDA et Qualite de Detection", "Visualisations conservees dans le rapport.")
    add_picture(slide, "race_gender_heatmap.png", Inches(0.65), Inches(1.5), width=Inches(5.6))
    add_picture(slide, "detection_rate_by_race.png", Inches(6.5), Inches(1.5), width=Inches(6.0))
    add_bullets(
        slide,
        Inches(0.8),
        Inches(5.1),
        Inches(12.0),
        Inches(1.5),
        [
            f"La detection FaceMesh varie selon les groupes et s'etend ici de {detection_min:.1f}% a {detection_max:.1f}%.",
            "Cet indicateur doit rester suivi car il influence directement la representativite de l'etape de modelisation.",
            f"Le pipeline n'a trouve aucune image manquante dans l'echantillon et seulement {metrics['duplicates']} doublons exacts.",
        ],
        font_size=15,
    )

    slide = add_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "5. Data Preparation", "Nettoyage, construction des variables et normalisation.")
    add_table(
        slide,
        Inches(0.7),
        Inches(1.45),
        Inches(6.1),
        Inches(4.5),
        ["Etape", "Resultat"],
        [
            ["Controle CSV", "Aucune valeur manquante dans les labels"],
            ["Verification fichiers", "0 image manquante sur l'echantillon"],
            ["Doublons exacts", str(metrics["duplicates"])],
            ["Extraction FaceMesh", f"{metrics['detected_faces']} visages detectes"],
            ["Valeurs manquantes features", "0 cellule manquante"],
            ["Outliers", f"{metrics['outliers_removed']} lignes retirees via z-score > 3"],
            ["Standardisation", f"{metrics['modeling_rows']} lignes conservees"],
        ],
    )
    add_bullets(
        slide,
        Inches(7.2),
        Inches(1.55),
        Inches(5.3),
        Inches(2.8),
        [
            "10 ratios geometriques normalises",
            "Features robustes a l'echelle absolue",
            "Aucun encodage categoriel necessaire pour le clustering",
            "Variables standardisees avant ACP et modelisation",
        ],
        font_size=16,
    )
    add_picture(slide, "feature_distributions.png", Inches(7.0), Inches(3.55), width=Inches(5.55))

    slide = add_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "6. ACP (PCA)", "Reduction de dimension et interpretation des composantes.")
    add_picture(slide, "pca_scree_plot.png", Inches(0.7), Inches(1.55), width=Inches(5.25))
    add_picture(slide, "pca_loadings_heatmap.png", Inches(6.35), Inches(1.45), width=Inches(5.9))
    pca_bullets = [
        f"ACP retenue : {metrics['pca_components']} composantes pour {metrics['pca_variance'] * 100:.2f}% de variance expliquee.",
    ]
    for component, names in metrics["top_features"]:
        pca_bullets.append(f"{component} : {', '.join(names)}")
    add_bullets(slide, Inches(0.8), Inches(5.1), Inches(12.0), Inches(1.6), pca_bullets, font_size=14)

    slide = add_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "7. Modeling", "Comparaison des modeles et justification du choix final.")
    add_picture(slide, "model_comparison.png", Inches(0.7), Inches(1.55), width=Inches(6.0))
    add_table(
        slide,
        Inches(7.0),
        Inches(1.55),
        Inches(5.55),
        Inches(2.4),
        ["Modele", "k", "Silhouette", "Calinski", "Davies"],
        [
            [
                row["model_label"],
                str(int(row["k"])),
                f"{row['silhouette']:.3f}",
                f"{row['calinski_harabasz']:.1f}",
                f"{row['davies_bouldin']:.3f}",
            ]
            for _, row in metrics["best_models"].iterrows()
        ],
    )
    add_bullets(
        slide,
        Inches(7.05),
        Inches(4.25),
        Inches(5.4),
        Inches(2.0),
        [
            f"Meilleur score observe : {best_model['model_label']} avec k={int(best_model['k'])}.",
            "Gaussian Mixture capture mieux des formes de clusters non strictement spherique.",
            "KMeans reste valorise dans le rapport pour sa coherence avec le pipeline voix precedent.",
        ],
        font_size=15,
    )

    slide = add_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "8. Interpretation des Clusters", "Lecture des profils issus du modele final.")
    add_picture(slide, "cluster_scatter.png", Inches(0.7), Inches(1.5), width=Inches(5.4))
    add_picture(slide, "cluster_profile_heatmap.png", Inches(6.35), Inches(1.5), width=Inches(6.0))
    add_table(
        slide,
        Inches(0.9),
        Inches(5.1),
        Inches(11.8),
        Inches(1.4),
        ["Cluster", "Part (%)", "Traits dominants", "Trait en retrait"],
        metrics["cluster_profiles"].values.tolist(),
    )

    slide = add_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "9. Conclusion", "Synthese pour la soutenance ou le rendu.")
    add_bullets(
        slide,
        Inches(0.85),
        Inches(1.55),
        Inches(11.9),
        Inches(3.0),
        [
            "Les sections exigees sont toutes couvertes : BO/DSO, EDA, preparation, ACP et modelisation.",
            "La presentation repose sur des mesures reellement executees et non sur un texte purement theorique.",
            "Le principal point d'attention reste l'amelioration du taux de detection FaceMesh.",
            "Pour une coherence metodologique maximale avec le module voix, KMeans peut etre retenu comme variante defendable.",
        ],
        font_size=18,
    )
    add_card(slide, Inches(1.05), Inches(5.0), Inches(2.3), Inches(1.2), "Detection", f"{metrics['detection_rate'] * 100:.2f}%", NAVY)
    add_card(slide, Inches(3.7), Inches(5.0), Inches(2.3), Inches(1.2), "ACP", f"{metrics['pca_variance'] * 100:.2f}%", TEAL)
    add_card(slide, Inches(6.35), Inches(5.0), Inches(2.8), Inches(1.2), "Modele final", f"{best_model['model_label']} / k={int(best_model['k'])}", CORAL)
    add_card(slide, Inches(9.55), Inches(5.0), Inches(2.1), Inches(1.2), "Option voix", f"KMeans / {kmeans_model['silhouette']:.3f}", OLIVE)

    prs.save(PPTX_PATH)
    return PPTX_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(f"Presentation generated at: {path}")
